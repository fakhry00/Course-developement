"""
Export utilities for converting content to various formats
"""

import markdown
from pathlib import Path
from typing import Optional
from docx import Document
from docx.shared import Inches
from pptx import Presentation
from pptx.util import Inches, Pt
from reportlab.lib import colors
from reportlab.lib.pagesizes import letter, A4
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, PageBreak
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import inch
import re

class ExportTools:
    """Utility class for exporting content to various formats"""
    
    def __init__(self):
        self.styles = getSampleStyleSheet()
    
    def markdown_to_pdf(self, markdown_content: str, output_path: Path):
        """Convert markdown content to PDF"""
        
        try:
            # Create PDF document
            doc = SimpleDocTemplate(
                str(output_path),
                pagesize=A4,
                rightMargin=72,
                leftMargin=72,
                topMargin=72,
                bottomMargin=18
            )
            
            # Parse markdown and create story
            story = self._markdown_to_reportlab_story(markdown_content)
            
            # Build PDF
            doc.build(story)
            
        except Exception as e:
            raise Exception(f"Error creating PDF: {str(e)}")
    
    def markdown_to_docx(self, markdown_content: str, output_path: Path):
        """Convert markdown content to Word document"""
        
        try:
            doc = Document()
            
            # Parse markdown content
            lines = markdown_content.split('\n')
            
            for line in lines:
                line = line.strip()
                
                if not line:
                    # Add space for empty lines
                    doc.add_paragraph()
                    continue
                
                # Handle headers
                if line.startswith('# '):
                    doc.add_heading(line[2:], level=1)
                elif line.startswith('## '):
                    doc.add_heading(line[3:], level=2)
                elif line.startswith('### '):
                    doc.add_heading(line[4:], level=3)
                elif line.startswith('#### '):
                    doc.add_heading(line[5:], level=4)
                # Handle bullet points
                elif line.startswith('- ') or line.startswith('* '):
                    doc.add_paragraph(line[2:], style='List Bullet')
                # Handle numbered lists
                elif re.match(r'^\d+\. ', line):
                    doc.add_paragraph(re.sub(r'^\d+\. ', '', line), style='List Number')
                # Regular paragraphs
                else:
                    doc.add_paragraph(line)
            
            doc.save(str(output_path))
            
        except Exception as e:
            raise Exception(f"Error creating Word document: {str(e)}")
    
    def markdown_to_pptx(self, markdown_content: str, output_path: Path):
        """Convert markdown content to PowerPoint presentation"""
        
        try:
            prs = Presentation()
            
            # Split content by slide breaks (---)
            slides_content = markdown_content.split('---')
            
            for slide_content in slides_content:
                slide_content = slide_content.strip()
                if not slide_content:
                    continue
                
                # Parse slide content
                lines = slide_content.split('\n')
                
                # First line as title, rest as content
                title = lines[0].strip('#').strip() if lines else "Slide"
                content_lines = [line for line in lines[1:] if line.strip()]
                
                # Add slide
                slide_layout = prs.slide_layouts[1]  # Title and Content layout
                slide = prs.slides.add_slide(slide_layout)
                
                # Set title
                slide.shapes.title.text = title
                
                # Add content
                if content_lines:
                    content_shape = slide.shapes.placeholders[1]
                    text_frame = content_shape.text_frame
                    text_frame.text = content_lines[0]
                    
                    # Add additional paragraphs
                    for line in content_lines[1:]:
                        if line.strip():
                            p = text_frame.add_paragraph()
                            p.text = line.strip('- *').strip()
            
            prs.save(str(output_path))
            
        except Exception as e:
            raise Exception(f"Error creating PowerPoint: {str(e)}")
    
    def _markdown_to_reportlab_story(self, markdown_content: str):
        """Convert markdown to ReportLab story elements"""
        
        story = []
        lines = markdown_content.split('\n')
        
        for line in lines:
            line = line.strip()
            
            if not line:
                story.append(Spacer(1, 12))
                continue
            
            # Handle headers
            if line.startswith('# '):
                story.append(Paragraph(line[2:], self.styles['Heading1']))
                story.append(Spacer(1, 12))
            elif line.startswith('## '):
                story.append(Paragraph(line[3:], self.styles['Heading2']))
                story.append(Spacer(1, 12))
            elif line.startswith('### '):
                story.append(Paragraph(line[4:], self.styles['Heading3']))
                story.append(Spacer(1, 6))
            elif line.startswith('#### '):
                story.append(Paragraph(line[5:], self.styles['Heading4']))
                story.append(Spacer(1, 6))
            # Handle bullet points
            elif line.startswith('- ') or line.startswith('* '):
                story.append(Paragraph(line[2:], self.styles['Bullet']))
            # Regular paragraphs
            else:
                story.append(Paragraph(line, self.styles['Normal']))
                story.append(Spacer(1, 6))
        
        return story
    
    def create_zip_archive(self, source_dir: Path, output_path: Path):
        """Create a ZIP archive from a directory"""
        
        import zipfile
        
        with zipfile.ZipFile(output_path, 'w', zipfile.ZIP_DEFLATED) as zipf:
            for file_path in source_dir.rglob('*'):
                if file_path.is_file():
                    arcname = file_path.relative_to(source_dir.parent)
                    zipf.write(file_path, arcname)
